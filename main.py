# 修改后的main.py（增强版）
import base64
import io
import json
import os
import uuid
import re
from typing import Any, Dict, Optional, Tuple, List, Set
from PIL import Image
import logging
import math

try:
    import requests
except ImportError:
    requests = None

# 可选 jsonschema 校验
try:
    from jsonschema import validate as _js_validate
except ImportError:  # 允许缺失
    _js_validate = None

try:
    from pptx import Presentation
    from pptx.util import Pt, Inches
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
    from pptx.chart.data import CategoryChartData, XyChartData, BubbleChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.enum.dml import MSO_LINE
except ImportError:
    Presentation = None

_IMAGE_CACHE: Dict[str, bytes] = {}
_THEME: Dict[str, Any] = {}  # 全局主题配置

# 安全/资源限制（可通过环境变量覆盖）
MAX_SLIDES = int(os.environ.get("JSON2PPT_MAX_SLIDES", "200"))
MAX_ELEMENTS_PER_SLIDE = int(os.environ.get("JSON2PPT_MAX_ELEMENTS", "600"))
MAX_REMOTE_IMAGE_BYTES = int(os.environ.get("JSON2PPT_MAX_IMG_BYTES", str(5 * 1024 * 1024)))  # 5MB
ALLOW_REMOTE_IMAGES = os.environ.get("JSON2PPT_ALLOW_REMOTE", "1") == "1"
ALLOW_FILE_IMAGES = os.environ.get("JSON2PPT_ALLOW_FILE", "1") == "1"
ALLOWED_REMOTE_DOMAINS: Optional[Set[str]] = None
if os.environ.get("JSON2PPT_REMOTE_DOMAINS"):
    ALLOWED_REMOTE_DOMAINS = {d.strip().lower() for d in os.environ["JSON2PPT_REMOTE_DOMAINS"].split(",") if d.strip()}
ASSET_ROOT = os.environ.get("JSON2PPT_ASSET_ROOT")  # 若设置则 file: 必须在此目录内

SCHEMA_CACHE: Optional[Dict[str, Any]] = None  # 缓存 schema


def hex_to_rgb(color: str):
    """转换颜色值，支持主题变量引用"""
    if not color:
        return (0, 0, 0)

    # 支持主题颜色变量 $primary, $secondary等
    if color.startswith('$'):
        var_name = color[1:]
        if _THEME and 'colors' in _THEME:
            color = _THEME['colors'].get(var_name, '#000000')

    c = color.strip().lstrip('#') if color else ''
    if len(c) == 3:
        c = ''.join(ch * 2 for ch in c)
    if len(c) == 6:
        return (int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16))
    elif len(c) == 8:  # 支持RGBA
        return (int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16), int(c[6:8], 16))
    return (0, 0, 0)


def unit_to_emu(value: float, total: int, unit: str) -> int:
    if unit == "percent":
        px = total * value / 100.0
    else:
        px = value
    return int(px * 9525)


def resolve_box(box: Dict[str, Any], slide_w: int, slide_h: int, default_unit: str):
    unit = box.get("unit", default_unit)
    x = unit_to_emu(box.get("x", 0), slide_w, unit)
    y = unit_to_emu(box.get("y", 0), slide_h, unit)
    w = unit_to_emu(box.get("w", slide_w), slide_w, unit)
    h = unit_to_emu(box.get("h", slide_h), slide_h, unit)
    return x, y, w, h


def apply_gradient(shape, gradient_cfg: Dict[str, Any]):
    """应用渐变效果"""
    if not gradient_cfg or not shape:
        return

    gradient_type = gradient_cfg.get("type", "linear")
    stops = gradient_cfg.get("stops", [])

    if not stops or len(stops) < 2:
        return

    try:
        shape.fill.gradient()

        if gradient_type == "linear":
            angle = gradient_cfg.get("angle", 0)
            shape.fill.gradient_angle = angle

        # 清除默认stops并添加新的
        shape.fill.gradient_stops.clear()
        for stop in stops:
            position = stop.get("position", 0) / 100.0  # 转换为0-1
            color = stop.get("color", "#000000")
            r, g, b = hex_to_rgb(color)[:3]
            gs = shape.fill.gradient_stops.add_stop(position)
            gs.color.rgb = RGBColor(r, g, b)
    except Exception as e:
        logging.warning(f"Failed to apply gradient: {e}")


def apply_shadow(shape, shadow_cfg: Dict[str, Any]):
    """应用阴影效果"""
    if not shadow_cfg or not shape:
        return

    try:
        shadow = shape.shadow
        shadow.inherit = False
        shadow.visible = True
        shadow.blur_radius = Pt(shadow_cfg.get("blur", 5))
        shadow.distance = Pt(math.sqrt(
            shadow_cfg.get("x", 2) ** 2 + shadow_cfg.get("y", 2) ** 2
        ))
        shadow.angle = math.degrees(math.atan2(
            shadow_cfg.get("y", 2), shadow_cfg.get("x", 2)
        ))

        color = shadow_cfg.get("color", "#000000")
        r, g, b = hex_to_rgb(color)[:3]
        shadow.color.rgb = RGBColor(r, g, b)
    except Exception as e:
        logging.warning(f"Failed to apply shadow: {e}")


def apply_rotation(shape, rotation: float):
    """应用旋转"""
    if rotation and shape:
        try:
            shape.rotation = rotation
        except:
            pass


def apply_opacity(shape, opacity: float):
    """应用透明度"""
    if opacity is not None and shape and 0 <= opacity <= 1:
        try:
            # python-pptx的透明度设置比较复杂，需要通过fill
            if hasattr(shape, 'fill'):
                shape.fill.transparency = 1.0 - opacity
        except:
            pass


def apply_hyperlink(shape, link_cfg: Dict[str, Any]):
    """应用超链接"""
    if not link_cfg or not shape:
        return

    try:
        if "url" in link_cfg:
            shape.click_action.hyperlink.address = link_cfg["url"]
        elif "slideId" in link_cfg:
            # 幻灯片内部链接需要特殊处理
            pass
    except:
        pass


def apply_border(shape, border_cfg: Dict[str, Any]):
    """应用边框"""
    if not border_cfg or not shape:
        return

    try:
        line = shape.line
        line.color.rgb = RGBColor(*hex_to_rgb(border_cfg.get("color", "#000000"))[:3])
        line.width = Pt(border_cfg.get("width", 1))

        style = border_cfg.get("style", "solid")
        if style == "dashed":
            line.dash_style = MSO_LINE.DASH
        elif style == "dotted":
            line.dash_style = MSO_LINE.DASH_DOT
        elif style == "double":
            # python-pptx不直接支持double
            pass
        else:
            line.dash_style = MSO_LINE.SOLID
    except Exception as e:
        logging.warning(f"Failed to apply border: {e}")


def ppt_color(fill_obj, color_hex: str):
    if not color_hex:
        return
    r, g, b = hex_to_rgb(color_hex)[:3]
    fill_obj.solid()
    fill_obj.fore_color.rgb = RGBColor(r, g, b)


def map_align(a: str):
    if a == "center":
        return PP_ALIGN.CENTER
    if a == "right":
        return PP_ALIGN.RIGHT
    if a == "justify":
        return PP_ALIGN.JUSTIFY
    return PP_ALIGN.LEFT


def map_valign(v: str):
    if v == "middle":
        return MSO_ANCHOR.MIDDLE
    if v == "bottom":
        return MSO_ANCHOR.BOTTOM
    return MSO_ANCHOR.TOP


def merge_styles(base: Dict[str, Any], override: Dict[str, Any]) -> Dict[str, Any]:
    merged = dict(base or {})
    merged.update(override or {})
    return merged


def apply_run_style(run, style: Dict[str, Any]):
    if not style:
        return
    if "fontSize" in style:
        run.font.size = Pt(style["fontSize"])
    if style.get("bold") is not None:
        run.font.bold = bool(style.get("bold"))
    if style.get("italic") is not None:
        run.font.italic = bool(style.get("italic"))
    if style.get("underline") is not None:
        run.font.underline = bool(style.get("underline"))
    if style.get("strikethrough") is not None:
        # python-pptx不直接支持strikethrough
        pass
    if style.get("color"):
        r, g, b = hex_to_rgb(style["color"])[:3]
        run.font.color.rgb = RGBColor(r, g, b)
    if "fontFamily" in style:
        font_name = style["fontFamily"]
        # 支持主题字体
        if font_name.startswith('$'):
            if _THEME and 'fonts' in _THEME:
                if font_name == '$heading':
                    font_name = _THEME['fonts'].get('heading', 'Arial')
                elif font_name == '$body':
                    font_name = _THEME['fonts'].get('body', 'Arial')
        run.font.name = font_name


def apply_paragraph_style(paragraph, style: Dict[str, Any]):
    if not style:
        return
    if style.get("align"):
        paragraph.alignment = map_align(style.get("align"))
    if "level" in style:
        try:
            paragraph.level = int(style["level"])
        except (ValueError, TypeError):
            pass
    if "lineHeight" in style:
        paragraph.line_spacing = style["lineHeight"]

    if paragraph.runs:
        run = paragraph.runs[0]
    else:
        run = paragraph.add_run()
    apply_run_style(run, style)


def apply_cell_style(cell, style: Dict[str, Any]):
    if not style:
        return
    if style.get("fill"):
        ppt_color(cell.fill, style["fill"])
    if style.get("gradient"):
        apply_gradient(cell, style["gradient"])

    text_frame = cell.text_frame
    if style.get("valign"):
        text_frame.vertical_anchor = map_valign(style.get("valign"))

    for paragraph in text_frame.paragraphs:
        apply_paragraph_style(paragraph, style)


def map_chart_type(t: str):
    mapping = {
        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "barStacked": XL_CHART_TYPE.COLUMN_STACKED,
        "barStacked100": XL_CHART_TYPE.COLUMN_STACKED_100,
        "barHorizontal": XL_CHART_TYPE.BAR_CLUSTERED,
        "line": XL_CHART_TYPE.LINE_MARKERS,
        "lineSmooth": XL_CHART_TYPE.LINE,
        "pie": XL_CHART_TYPE.PIE,
        "doughnut": XL_CHART_TYPE.DOUGHNUT,
        "area": XL_CHART_TYPE.AREA,
        "areaStacked": XL_CHART_TYPE.AREA_STACKED,
        "scatter": XL_CHART_TYPE.XY_SCATTER_LINES,
        "bubble": XL_CHART_TYPE.BUBBLE,
        "radar": XL_CHART_TYPE.RADAR
    }
    return mapping.get(t, XL_CHART_TYPE.COLUMN_CLUSTERED)


def pick_shape(shape_type: str):
    """扩展形状类型映射"""
    mapping = {
        "rect": MSO_SHAPE.RECTANGLE,
        "roundRect": MSO_SHAPE.ROUNDED_RECTANGLE,
        "ellipse": MSO_SHAPE.OVAL,
        "circle": MSO_SHAPE.OVAL,
        "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
        "rightTriangle": MSO_SHAPE.RIGHT_TRIANGLE,
        "pentagon": MSO_SHAPE.PENTAGON,
        "hexagon": MSO_SHAPE.HEXAGON,
        "octagon": MSO_SHAPE.OCTAGON,
        "star": MSO_SHAPE.STAR_5_POINT,
        "star5": MSO_SHAPE.STAR_5_POINT,
        "star6": MSO_SHAPE.STAR_6_POINT,
        "star8": MSO_SHAPE.STAR_8_POINT,
        "arrow": MSO_SHAPE.RIGHT_ARROW,
        "arrowLeft": MSO_SHAPE.LEFT_ARROW,
        "arrowRight": MSO_SHAPE.RIGHT_ARROW,
        "arrowUp": MSO_SHAPE.UP_ARROW,
        "arrowDown": MSO_SHAPE.DOWN_ARROW,
        "arrowBoth": MSO_SHAPE.LEFT_RIGHT_ARROW,
        "chevron": MSO_SHAPE.CHEVRON,
        "callout": MSO_SHAPE.ROUNDED_RECTANGULAR_CALLOUT,
        "cloud": MSO_SHAPE.CLOUD,
        "heart": MSO_SHAPE.HEART,
        "diamond": MSO_SHAPE.DIAMOND,
        "parallelogram": MSO_SHAPE.PARALLELOGRAM,
        "trapezoid": MSO_SHAPE.TRAPEZOID,
        "cross": MSO_SHAPE.CROSS,
        "plus": MSO_SHAPE.MATH_PLUS
    }
    return mapping.get(shape_type, MSO_SHAPE.RECTANGLE)


def _domain_allowed(url: str) -> bool:
    if not ALLOWED_REMOTE_DOMAINS:
        return True
    try:
        m = re.match(r'^https?://([^/]+)', url, re.IGNORECASE)
        if not m:
            return False
        host = m.group(1).lower()
        # 简单 host 匹配或子域匹配
        return any(host == d or host.endswith('.' + d) for d in ALLOWED_REMOTE_DOMAINS)
    except Exception:
        return False


def _check_file_path(path: str) -> bool:
    if not ALLOW_FILE_IMAGES:
        return False
    if ASSET_ROOT:
        try:
            real_root = os.path.realpath(ASSET_ROOT)
            real_path = os.path.realpath(path)
            return real_path.startswith(real_root)
        except Exception:
            return False
    return True


def get_image_bytes(source: str, logger=None):
    if not source:
        return None
    cached = _IMAGE_CACHE.get(source)
    if cached is not None:
        return cached
    if source.startswith("base64:"):
        b64 = source[len("base64:"):]
        try:
            # 自动补齐缺失的 '=' padding
            missing = len(b64) % 4
            if missing:
                b64 += '=' * (4 - missing)
            content = base64.b64decode(b64)
            _IMAGE_CACHE[source] = content
            return content
        except Exception as e:
            # 再尝试忽略非 base64 字符
            try:
                filtered = ''.join(ch for ch in b64 if ch.isalnum() or ch in '+/=' )
                content = base64.b64decode(filtered)
                _IMAGE_CACHE[source] = content
                return content
            except Exception:
                logger and logger.warning(f"base64 decode failed: {e}")
                return None
    elif source.startswith("url:"):
        if not ALLOW_REMOTE_IMAGES:
            logger and logger.warning("remote image disabled by config")
            return None
        if requests is None:
            logger and logger.warning("requests not installed, cannot fetch remote image")
            return None
        url = source[len("url:"):]
        if not url.lower().startswith(("http://", "https://")):
            logger and logger.warning("unsupported url scheme")
            return None
        if not _domain_allowed(url):
            logger and logger.warning(f"domain not allowed: {url}")
            return None
        try:
            resp = requests.get(url, timeout=(3, 7), stream=True)
            resp.raise_for_status()
            content_type = resp.headers.get("Content-Type", "").lower()
            if not any(ct in content_type for ct in ["image/", "application/octet-stream"]):
                logger and logger.warning(f"unexpected content-type {content_type}")
                return None
            data = b''
            total = 0
            for chunk in resp.iter_content(65536):
                if not chunk:
                    continue
                data += chunk
                total += len(chunk)
                if total > MAX_REMOTE_IMAGE_BYTES:
                    logger and logger.warning("remote image exceeds size limit")
                    return None
            _IMAGE_CACHE[source] = data
            return data
        except Exception as e:
            logger and logger.warning(f"download image failed: {e}")
        return None
    elif source.startswith("file:"):
        path = source[len("file:"):]
        if not _check_file_path(path):
            logger and logger.warning("file path not allowed")
            return None
        try:
            with open(path, "rb") as f:
                content = f.read()
                if len(content) > MAX_REMOTE_IMAGE_BYTES:
                    logger and logger.warning("local image exceeds size limit")
                    return None
                _IMAGE_CACHE[source] = content
                return content
        except Exception as e:
            logger and logger.warning(f"read file image failed: {e}")
            return None
    return None


def add_background(slide, bg_cfg: Dict[str, Any], slide_w_emu: int, slide_h_emu: int, logger=None):
    if not bg_cfg:
        return

    # 背景色
    color = bg_cfg.get("color")
    if color:
        ppt_color(slide.background.fill, color)

    # 背景渐变
    gradient = bg_cfg.get("gradient")
    if gradient:
        apply_gradient(slide.background, gradient)

    # 背景图片
    img = bg_cfg.get("image")
    if img and isinstance(img, dict):
        content = get_image_bytes(img.get("src"), logger)
        if content:
            from io import BytesIO
            pic = slide.shapes.add_picture(BytesIO(content), 0, 0, width=slide_w_emu, height=slide_h_emu)

            # 应用图片效果
            if "opacity" in img:
                apply_opacity(pic, img["opacity"])

            # 将图片移到最底层
            slide.shapes._spTree.remove(pic._element)
            slide.shapes._spTree.insert(2, pic._element)


def apply_base_element_properties(shape, elem: Dict[str, Any]):
    """应用基础元素属性"""
    # 旋转
    if "rotation" in elem:
        apply_rotation(shape, elem["rotation"])

    # 透明度
    if "opacity" in elem:
        apply_opacity(shape, elem["opacity"])

    # 阴影
    if "shadow" in elem:
        apply_shadow(shape, elem["shadow"])

    # 超链接
    if "link" in elem:
        apply_hyperlink(shape, elem["link"])


def add_text(slide, elem, slide_w, slide_h, default_unit):
    box = elem.get("box", {})
    x, y, w, h = resolve_box(box, slide_w, slide_h, default_unit)
    shape = slide.shapes.add_textbox(x, y, w, h)

    # 应用基础属性
    apply_base_element_properties(shape, elem)

    # 背景和边框
    if elem.get("fill"):
        ppt_color(shape.fill, elem["fill"])
    if elem.get("border"):
        apply_border(shape, elem["border"])

    # 内边距
    padding = elem.get("padding")
    if padding is not None:
        tf = shape.text_frame
        if isinstance(padding, dict):
            tf.margin_top = Pt(padding.get("top", 0))
            tf.margin_bottom = Pt(padding.get("bottom", 0))
            tf.margin_left = Pt(padding.get("left", 0))
            tf.margin_right = Pt(padding.get("right", 0))
        else:
            tf.margin_top = tf.margin_bottom = Pt(padding)
            tf.margin_left = tf.margin_right = Pt(padding)

    tf = shape.text_frame
    tf.clear()

    # 垂直对齐
    style = elem.get("style", {})
    if style.get("valign"):
        tf.vertical_anchor = map_valign(style["valign"])

    base_style = elem.get("style", {})
    paragraphs = elem.get("paragraphs")
    if paragraphs:
        for idx, para_cfg in enumerate(paragraphs):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            text = para_cfg.get("text", "")

            # 列表处理
            list_type = para_cfg.get("listType")
            if list_type == "number":
                number_index = para_cfg.get("number", idx + 1)
                prefix = para_cfg.get("numberPrefix", "")
                if prefix:
                    text = f"{prefix}{text}"
                else:
                    text = f"{number_index}. {text}"
            elif list_type == "bullet" or para_cfg.get("bullet"):
                bullet_char = para_cfg.get("bulletChar", "•")
                text = f"{bullet_char} {text}"

            # 缩进
            if "indent" in para_cfg:
                p.level = int(para_cfg["indent"])

            p.text = text
            final_style = merge_styles(base_style, para_cfg.get("style"))
            apply_paragraph_style(p, final_style)
    else:
        p = tf.paragraphs[0]
        p.text = elem.get("text", "")
        apply_paragraph_style(p, merge_styles({}, base_style))


def add_image(slide, elem, slide_w, slide_h, default_unit, logger=None):
    box = elem.get("box", {})
    x, y, w, h = resolve_box(box, slide_w, slide_h, default_unit)
    content = get_image_bytes(elem.get("source"), logger)
    if not content:
        return

    from io import BytesIO
    pic = slide.shapes.add_picture(BytesIO(content), x, y, width=w, height=h)

    # 应用基础属性
    apply_base_element_properties(pic, elem)

    # 边框
    if elem.get("border"):
        apply_border(pic, elem["border"])

    # 图片裁剪（python-pptx支持有限）
    crop = elem.get("crop")
    if crop:
        # 需要更复杂的实现
        pass


def add_shape(slide, elem, slide_w, slide_h, default_unit):
    box = elem.get("box", {})
    x, y, w, h = resolve_box(box, slide_w, slide_h, default_unit)
    shape_type = elem.get("shapeType", "rect")

    # 对于圆形，确保宽高相等
    if shape_type == "circle":
        size = min(w, h)
        w = h = size

    shape = slide.shapes.add_shape(pick_shape(shape_type), x, y, w, h)

    # 应用基础属性
    apply_base_element_properties(shape, elem)

    # 填充
    fill_color = elem.get("fill")
    if fill_color:
        ppt_color(shape.fill, fill_color)

    # 渐变
    gradient = elem.get("gradient")
    if gradient:
        apply_gradient(shape, gradient)

    # 边框
    if elem.get("border"):
        apply_border(shape, elem["border"])


def add_line(slide, elem, slide_w, slide_h, default_unit):
    """添加线条元素"""
    points = elem.get("points", [])
    if len(points) < 2:
        return

    # 简化实现：只画第一条线段
    p1 = points[0]
    p2 = points[1] if len(points) > 1 else points[0]

    x1 = unit_to_emu(p1.get("x", 0), slide_w, "px")
    y1 = unit_to_emu(p1.get("y", 0), slide_h, "px")
    x2 = unit_to_emu(p2.get("x", 100), slide_w, "px")
    y2 = unit_to_emu(p2.get("y", 100), slide_h, "px")

    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        x1, y1, x2, y2
    )

    # 线条样式
    line = connector.line
    if elem.get("stroke"):
        line.color.rgb = RGBColor(*hex_to_rgb(elem["stroke"])[:3])
    if elem.get("strokeWidth"):
        line.width = Pt(elem["strokeWidth"])

    style = elem.get("strokeStyle", "solid")
    if style == "dashed":
        line.dash_style = MSO_LINE.DASH
    elif style == "dotted":
        line.dash_style = MSO_LINE.DASH_DOT

    # 应用基础属性
    apply_base_element_properties(connector, elem)


def add_icon(slide, elem, slide_w, slide_h, default_unit):
    """添加图标元素（简化实现：使用形状代替）"""
    box = elem.get("box", {})
    x, y, w, h = resolve_box(box, slide_w, slide_h, default_unit)

    # 使用圆形作为图标占位符
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)

    # 颜色
    color = elem.get("color", "#000000")
    if color:
        ppt_color(shape.fill, color)

    # 渐变
    gradient = elem.get("gradient")
    if gradient:
        apply_gradient(shape, gradient)

    # 应用基础属性
    apply_base_element_properties(shape, elem)


def add_group(slide, elem, slide_w, slide_h, default_unit, logger=None):
    """添加组元素（递归处理子元素）"""
    elements = elem.get("elements", [])

    # python-pptx不直接支持组，所以我们依次添加元素
    for sub_elem in elements:
        t = sub_elem.get("type")
        try:
            if t == "text":
                add_text(slide, sub_elem, slide_w, slide_h, default_unit)
            elif t == "image":
                add_image(slide, sub_elem, slide_w, slide_h, default_unit, logger)
            elif t == "shape":
                add_shape(slide, sub_elem, slide_w, slide_h, default_unit)
            elif t == "chart":
                add_chart(slide, sub_elem, slide_w, slide_h, default_unit)
            elif t == "table":
                add_table(slide, sub_elem, slide_w, slide_h, default_unit)
            elif t == "line":
                add_line(slide, sub_elem, slide_w, slide_h, default_unit)
            elif t == "icon":
                add_icon(slide, sub_elem, slide_w, slide_h, default_unit)
            elif t == "group":
                add_group(slide, sub_elem, slide_w, slide_h, default_unit, logger)
        except Exception as e:
            logger and logger.warning(f"Group element failed: {e}")


def add_table(slide, elem, slide_w, slide_h, default_unit):
    box = elem.get("box", {})
    x, y, w, h = resolve_box(box, slide_w, slide_h, default_unit)
    table_cfg = elem.get("table", {})
    header = table_cfg.get("header") or []
    rows = table_cfg.get("rows") or []
    if not header and not rows:
        return
    column_count = len(header) if header else (len(rows[0]) if rows else 0)
    if column_count == 0:
        return
    total_rows = len(rows) + (1 if header else 0)
    shape = slide.shapes.add_table(total_rows, column_count, x, y, w, h)
    table = shape.table

    # 应用基础属性
    apply_base_element_properties(shape, elem)

    # 表格边框
    if table_cfg.get("border"):
        # python-pptx表格边框设置比较复杂
        pass

    column_widths = table_cfg.get("columnWidths") or []
    width_unit = table_cfg.get("columnWidthsUnit", default_unit)
    for idx, width in enumerate(column_widths):
        if idx < len(table.columns):
            table.columns[idx].width = unit_to_emu(width, slide_w, width_unit)

    style_cfg = table_cfg.get("style") or {}
    base_cell_style = {k: v for k, v in style_cfg.items() if k not in {"header", "body", "footer"}}
    body_style = merge_styles(base_cell_style, style_cfg.get("body"))
    header_style = merge_styles(body_style, style_cfg.get("header"))

    def populate_row(row_idx, values, style):
        for col_idx in range(column_count):
            cell = table.cell(row_idx, col_idx)
            text = values[col_idx] if col_idx < len(values) else ""
            cell.text = "" if text is None else str(text)
            apply_cell_style(cell, style)

    current_row = 0
    if header:
        populate_row(current_row, header, header_style)
        current_row += 1
    for row_values in rows:
        populate_row(current_row, row_values, body_style)
        current_row += 1

    if table_cfg.get("bandedRows"):
        table.has_banded_rows = True
    if table_cfg.get("bandedColumns"):
        table.has_banded_columns = True


def add_chart(slide, elem, slide_w, slide_h, default_unit, logger=None):
    """添加图表，支持常见分类/XY/气泡类型，增强健壮性防止空图。"""
    box = elem.get("box", {})
    x, y, w, h = resolve_box(box, slide_w, slide_h, default_unit)
    chart_type_str = elem.get("chartType", "bar")
    chart_type = map_chart_type(chart_type_str)
    data_cfg = elem.get("data", {}) or {}
    raw_categories = data_cfg.get("categories") or []
    series_cfg = data_cfg.get("series") or []

    # 预处理：去除 None / 非数字（保留可转 float）
    def sanitize_values(vals):
        safe = []
        for v in (vals or []):
            if v is None:
                safe.append(None)
            else:
                try:
                    safe.append(float(v))
                except Exception:
                    # 非数字转 None
                    safe.append(None)
        return safe

    # 分类型图表集合
    category_types = {
        'COLUMN', 'BAR', 'LINE', 'AREA', 'PIE', 'DOUGHNUT', 'RADAR', 'WATERFALL',
        'COLUMN_STACKED', 'COLUMN_STACKED_100', 'BAR_STACKED', 'BAR_STACKED_100'
    }

    # XY / 散点 / 气泡类型检测（通过映射的枚举名称）
    chart_type_name = getattr(chart_type, 'name', str(chart_type))
    is_xy = 'SCATTER' in chart_type_name or 'XY_SCATTER' in chart_type_name
    is_bubble = 'BUBBLE' in chart_type_name

    try:
        if is_bubble:
            # 气泡：每个点需要 (x, y, size)
            bdata = BubbleChartData()
            # 气泡图无 categories，按 series 顺序
            if not series_cfg:
                # 强制一个空系列避免 PPTX 生成空框
                series_cfg = [{"name": "Series", "values": [(0, 0, 1)]}]
            for s in series_cfg:
                name = s.get("name", "Series")
                values = s.get("values", [])
                # 允许 values 为 [ {x:.., y:.., size:..}, ... ] 或 [[x,y,size], ...]
                cleaned = []
                for item in values:
                    if isinstance(item, dict):
                        x_v = item.get('x', 0)
                        y_v = item.get('y', 0)
                        sz = item.get('size', 1)
                    elif isinstance(item, (list, tuple)) and len(item) >= 3:
                        x_v, y_v, sz = item[0], item[1], item[2]
                    else:
                        continue
                    try:
                        x_v = float(x_v); y_v = float(y_v); sz = float(sz) if sz else 1.0
                    except Exception:
                        continue
                    cleaned.append((x_v, y_v, sz))
                if not cleaned:
                    cleaned = [(0.0, 0.0, 1.0)]
                serie = bdata.add_series(name)
                for x_v, y_v, sz in cleaned:
                    serie.add_data_point(x_v, y_v, sz)
            data_obj = bdata
        elif is_xy:
            xydata = XyChartData()
            if not series_cfg:
                series_cfg = [{"name": "Series", "values": [(0, 0)]}]
            for s in series_cfg:
                name = s.get("name", "Series")
                values = s.get("values", [])
                serie = xydata.add_series(name)
                for pair in values:
                    if isinstance(pair, (list, tuple)) and len(pair) >= 2:
                        try:
                            x_v = float(pair[0]); y_v = float(pair[1])
                        except Exception:
                            continue
                        serie.add_data_point(x_v, y_v)
                    elif isinstance(pair, dict) and 'x' in pair and 'y' in pair:
                        try:
                            serie.add_data_point(float(pair['x']), float(pair['y']))
                        except Exception:
                            continue
                if not list(serie._points):  # 没有有效点
                    serie.add_data_point(0.0, 0.0)
            data_obj = xydata
        else:
            # 分类图
            cdata = CategoryChartData()
            # 空 categories 但有 series，以 series 长度推一个序号分类
            if not raw_categories and series_cfg:
                longest = max((len(s.get('values', [])) for s in series_cfg), default=0)
                raw_categories = [f'C{i+1}' for i in range(longest)]
            cdata.categories = [str(c) for c in raw_categories]
            cat_len = len(cdata.categories)
            for s in series_cfg:
                name = s.get("name", "Series")
                values = sanitize_values(s.get("values", []))
                # 修剪 / 填充长度
                if cat_len and len(values) != cat_len:
                    if len(values) > cat_len:
                        values = values[:cat_len]
                    else:
                        values = values + [None] * (cat_len - len(values))
                cdata.add_series(name, values)
            # 如果没有系列生成一个空系列确保渲染框
            if not series_cfg:
                cdata.add_series("Series", [0])
            data_obj = cdata

        chart_shape = slide.shapes.add_chart(chart_type, x, y, w, h, data_obj)
        chart = chart_shape.chart
    except Exception as e:
        if logger:
            logger.warning(f"add_chart failed ({chart_type_str}): {e}")
        # 回退到一个空的默认柱状图，避免整个生成失败
        fallback = CategoryChartData(); fallback.categories = ['C1']; fallback.add_series('Series', [0])
        chart_shape = slide.shapes.add_chart(map_chart_type('bar'), x, y, w, h, fallback)
        chart = chart_shape.chart

    # 应用基础属性
    apply_base_element_properties(chart_shape, elem)

    title = elem.get("title")
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = title

    options = elem.get("chartOptions", {})

    # 图例
    legend_cfg = options.get("legend")
    if legend_cfg is not None:
        if isinstance(legend_cfg, bool):
            chart.has_legend = legend_cfg
        elif isinstance(legend_cfg, dict):
            chart.has_legend = legend_cfg.get("show", True)
            if chart.has_legend and "position" in legend_cfg:
                # 设置图例位置
                pass

    # 数据标签
    if "dataLabels" in options:
        data_labels_cfg = options.get("dataLabels")
        if isinstance(data_labels_cfg, bool):
            enabled = data_labels_cfg
        else:
            enabled = data_labels_cfg.get("enabled", True)
        try:
            if hasattr(chart, 'plots'):
                for plot in chart.plots:
                    plot.has_data_labels = enabled
        except:
            pass

    # 网格线
    grid_cfg = options.get("grid")
    if grid_cfg and hasattr(chart, "value_axis"):
        try:
            chart.value_axis.has_major_gridlines = grid_cfg.get("show", True)
        except:
            pass

    # 值轴
    if hasattr(chart, "value_axis") and chart.value_axis and isinstance(options.get("valueAxis"), dict):
        try:
            axis_cfg = options["valueAxis"]
            value_axis = chart.value_axis
            if "minimum" in axis_cfg:
                value_axis.minimum_scale = axis_cfg["minimum"]
            if "maximum" in axis_cfg:
                value_axis.maximum_scale = axis_cfg["maximum"]
            if "majorUnit" in axis_cfg:
                value_axis.major_unit = axis_cfg["majorUnit"]
            if "title" in axis_cfg:
                value_axis.has_title = True
                value_axis.axis_title.text_frame.text = axis_cfg["title"]
        except Exception:
            pass

    # 分类轴
    if hasattr(chart, "category_axis") and chart.category_axis and isinstance(options.get("categoryAxis"), dict):
        try:
            cat_cfg = options["categoryAxis"]
            category_axis = chart.category_axis
            if "tickLabelRotation" in cat_cfg:
                category_axis.tick_label_rotation = cat_cfg["tickLabelRotation"]
            if "title" in cat_cfg:
                category_axis.has_title = True
                category_axis.axis_title.text_frame.text = cat_cfg["title"]
        except Exception:
            pass

    # 设置系列颜色和样式
    colors = options.get("colors", [])
    for idx, series in enumerate(chart.series):
        color = None
        gradient = None

        # 从series配置获取颜色
        if idx < len(series_cfg):
            color = series_cfg[idx].get("color")
            gradient = series_cfg[idx].get("gradient")

        # 从options.colors获取颜色
        if not color and idx < len(colors):
            color = colors[idx]

        if color:
            try:
                r, g, b = hex_to_rgb(color)[:3]
                # PIE / DOUGHNUT 使用 points 着色
                if hasattr(series, 'format') and series.format and hasattr(series.format, 'fill'):
                    series.format.fill.solid()
                    series.format.fill.fore_color.rgb = RGBColor(r, g, b)
                # 雷达 / 线图等保证 markers/line
                if hasattr(series, 'marker') and series.marker and hasattr(series.marker, 'format'):
                    try:
                        series.marker.format.fill.solid(); series.marker.format.fill.fore_color.rgb = RGBColor(r, g, b)
                    except Exception:
                        pass
            except Exception:
                pass
        elif gradient:
            # 暂不实现图表渐变，保留结构
            pass


def add_video(slide, elem, slide_w, slide_h, default_unit, logger=None):
    """添加视频元素（占位符实现）"""
    box = elem.get("box", {})
    x, y, w, h = resolve_box(box, slide_w, slide_h, default_unit)

    # python-pptx的视频支持有限，使用占位符
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0, 0, 0)

    # 添加文字说明
    shape.text_frame.text = "📹 Video Placeholder"
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 应用基础属性
    apply_base_element_properties(shape, elem)

    if elem.get("border"):
        apply_border(shape, elem["border"])


def add_smartart(slide, elem, slide_w, slide_h, default_unit, logger=None):
    """添加SmartArt元素（简化实现）"""
    box = elem.get("box", {})
    x, y, w, h = resolve_box(box, slide_w, slide_h, default_unit)

    smartart_type = elem.get("smartArtType", "list")
    nodes = elem.get("nodes", [])

    # python-pptx不直接支持SmartArt，使用形状组合模拟
    # 这里简化为列表展示
    if smartart_type in ["list", "process"]:
        node_h = h / max(len(nodes), 1)
        for i, node in enumerate(nodes):
            node_y = y + i * node_h

            # 添加形状
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                x, node_y, w * 0.8, node_h * 0.8
            )

            # 设置颜色
            if node.get("color"):
                ppt_color(shape.fill, node["color"])
            else:
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(91, 155, 213)

            # 添加文本
            shape.text_frame.text = node.get("text", f"Item {i + 1}")
            shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 应用基础属性
    apply_base_element_properties(slide.shapes[-1] if slide.shapes else None, elem)


def apply_transition(slide, transition_cfg: Dict[str, Any]):
    """应用幻灯片过渡效果（python-pptx支持有限）"""
    if not transition_cfg:
        return

    # python-pptx不直接支持过渡效果
    # 需要通过XML操作来实现
    pass


def resolve_slide_layout(prs, layout_cfg, fallback):
    if layout_cfg is None:
        return fallback
    if isinstance(layout_cfg, int):
        try:
            return prs.slide_layouts[layout_cfg]
        except IndexError:
            return fallback
    if isinstance(layout_cfg, str):
        if layout_cfg.isdigit():
            return resolve_slide_layout(prs, int(layout_cfg), fallback)
        preset_map = {
            "title": 0,
            "titleandcontent": 1,
            "sectionheader": 2,
            "twocontent": 3,
            "comparison": 4,
            "titleonly": 5,
            "blank": 6
        }
        key = layout_cfg.replace(" ", "").lower()
        if key in preset_map:
            idx = preset_map[key]
            if idx < len(prs.slide_layouts):
                return prs.slide_layouts[idx]
        for layout in prs.slide_layouts:
            name = getattr(layout, "name", None)
            if name and name.replace(" ", "").lower() == key:
                return layout
    return fallback


def load_schema(logger=None):
    global SCHEMA_CACHE
    if SCHEMA_CACHE is not None:
        return SCHEMA_CACHE
    schema_path = os.path.join(os.path.dirname(__file__), 'ppt.schema.json')
    if not os.path.exists(schema_path):
        return None
    try:
        with open(schema_path, 'r', encoding='utf-8') as f:
            SCHEMA_CACHE = json.load(f)
            return SCHEMA_CACHE
    except Exception as e:
        logger and logger.warning(f"load schema failed: {e}")
        return None


def validate(meta, use_schema: bool = True, logger=None):
    if "ppt" not in meta:
        raise ValueError("missing ppt root")
    slides = meta["ppt"].get("slides")
    if not isinstance(slides, list):
        raise ValueError("ppt.slides must be list")
    if len(slides) > MAX_SLIDES:
        raise ValueError(f"slides exceed limit {len(slides)}/{MAX_SLIDES}")
    for i, s in enumerate(slides):
        elems = s.get("elements", [])
        if not isinstance(elems, list):
            raise ValueError(f"slide[{i}].elements must be list")
        if len(elems) > MAX_ELEMENTS_PER_SLIDE:
            raise ValueError(f"slide[{i}] elements exceed limit {len(elems)}/{MAX_ELEMENTS_PER_SLIDE}")

    if use_schema and _js_validate:
        schema = load_schema(logger)
        if schema:
            try:
                _js_validate(instance=meta, schema=schema)
            except Exception as e:
                raise ValueError(f"schema validation failed: {e}")
    return True


def build_single_slide(meta: Dict[str, Any], slide_index: int, logger=None) -> Optional[Presentation]:
    """构建单个幻灯片的PPT，用于预览"""
    if Presentation is None:
        raise RuntimeError("python-pptx not installed")

    ppt_cfg = meta["ppt"]

    # 设置全局主题
    global _THEME
    _THEME = ppt_cfg.get("theme", {})

    slides_cfg = ppt_cfg.get("slides", [])

    if slide_index < 0 or slide_index >= len(slides_cfg):
        return None

    prs = Presentation()

    # 设置幻灯片尺寸
    size_cfg = ppt_cfg.get("size", {"width": 1280, "height": 720, "unit": "px"})
    sw = size_cfg.get("width", 1280)
    sh = size_cfg.get("height", 720)
    prs.slide_width = unit_to_emu(sw, sw, "px")
    prs.slide_height = unit_to_emu(sh, sh, "px")

    default_unit = ppt_cfg.get("defaultUnit", "px")
    default_layout = resolve_slide_layout(prs, ppt_cfg.get("defaultLayout"), prs.slide_layouts[6])

    # 只构建指定的幻灯片
    s = slides_cfg[slide_index]
    slide_layout = resolve_slide_layout(prs, s.get("layout"), default_layout)
    slide = prs.slides.add_slide(slide_layout)

    # 背景
    add_background(slide, s.get("background"), prs.slide_width, prs.slide_height, logger)

    # 过渡效果
    apply_transition(slide, s.get("transition"))

    # 处理元素（支持zIndex排序）
    elements = s.get("elements", [])
    sorted_elements = sorted(elements, key=lambda e: e.get("zIndex", 0))

    for elem in sorted_elements:
        t = elem.get("type")
        try:
            if t == "text":
                add_text(slide, elem, sw, sh, default_unit)
            elif t == "image":
                add_image(slide, elem, sw, sh, default_unit, logger)
            elif t == "chart":
                add_chart(slide, elem, sw, sh, default_unit)
            elif t == "shape":
                add_shape(slide, elem, sw, sh, default_unit)
            elif t == "table":
                add_table(slide, elem, sw, sh, default_unit)
            elif t == "line":
                add_line(slide, elem, sw, sh, default_unit)
            elif t == "icon":
                add_icon(slide, elem, sw, sh, default_unit)
            elif t == "group":
                add_group(slide, elem, sw, sh, default_unit, logger)
            elif t == "video":
                add_video(slide, elem, sw, sh, default_unit, logger)
            elif t == "smartArt":
                add_smartart(slide, elem, sw, sh, default_unit, logger)
        except Exception as e:
            logger and logger.warning(f"element failed id={elem.get('id')} err={e}")

    return prs


def build(meta, logger=None):
    if Presentation is None:
        raise RuntimeError("python-pptx not installed")

    prs = Presentation()
    ppt_cfg = meta["ppt"]

    # 设置全局主题
    global _THEME
    _THEME = ppt_cfg.get("theme", {})

    size_cfg = ppt_cfg.get("size", {"width": 1280, "height": 720, "unit": "px"})
    sw = size_cfg.get("width", 1280)
    sh = size_cfg.get("height", 720)
    prs.slide_width = unit_to_emu(sw, sw, "px")
    prs.slide_height = unit_to_emu(sh, sh, "px")

    default_unit = ppt_cfg.get("defaultUnit", "px")
    default_layout = resolve_slide_layout(prs, ppt_cfg.get("defaultLayout"), prs.slide_layouts[6])
    slides_cfg = ppt_cfg["slides"]

    for s in slides_cfg:
        slide_layout = resolve_slide_layout(prs, s.get("layout"), default_layout)
        slide = prs.slides.add_slide(slide_layout)

        # 背景
        add_background(slide, s.get("background"), prs.slide_width, prs.slide_height, logger)

        # 过渡效果
        apply_transition(slide, s.get("transition"))

        # 处理元素（支持zIndex排序）
        elements = s.get("elements", [])
        sorted_elements = sorted(elements, key=lambda e: e.get("zIndex", 0))

        for elem in sorted_elements:
            t = elem.get("type")
            try:
                if t == "text":
                    add_text(slide, elem, sw, sh, default_unit)
                elif t == "image":
                    add_image(slide, elem, sw, sh, default_unit, logger)
                elif t == "chart":
                    add_chart(slide, elem, sw, sh, default_unit)
                elif t == "shape":
                    add_shape(slide, elem, sw, sh, default_unit)
                elif t == "table":
                    add_table(slide, elem, sw, sh, default_unit)
                elif t == "line":
                    add_line(slide, elem, sw, sh, default_unit)
                elif t == "icon":
                    add_icon(slide, elem, sw, sh, default_unit)
                elif t == "group":
                    add_group(slide, elem, sw, sh, default_unit, logger)
                elif t == "video":
                    add_video(slide, elem, sw, sh, default_unit, logger)
                elif t == "smartArt":
                    add_smartart(slide, elem, sw, sh, default_unit, logger)
            except Exception as e:
                logger and logger.warning(f"element failed id={elem.get('id')} err={e}")

    return prs, len(slides_cfg)


def save_to_temp(prs) -> str:
    tmp_dir = "./"
    name = f"ppt_{uuid.uuid4().hex}.pptx"
    path = os.path.abspath(os.path.join(tmp_dir, name))
    prs.save(path)
    return path


def handler(args):
    logger = getattr(args, 'logger', None)
    meta_str = getattr(args.input, 'mets', None) or getattr(args.input, 'meta', None)
    if not meta_str:
        return {"message": "missing mets/meta", "page_count": 0, "file_path": ""}
    try:
        meta = json.loads(meta_str)
        validate(meta, use_schema=True, logger=logger)
        prs, count = build(meta, logger)
        file_path = save_to_temp(prs)
        return {
            "message": "OK",
            "page_count": count,
            "file_path": file_path,
            "limits": {
                "max_slides": MAX_SLIDES,
                "max_elements_per_slide": MAX_ELEMENTS_PER_SLIDE,
                "max_image_bytes": MAX_REMOTE_IMAGE_BYTES,
                "remote_images": ALLOW_REMOTE_IMAGES,
                "file_images": ALLOW_FILE_IMAGES
            }
        }
    except Exception as e:
        if logger:
            logger.error(f"failed: {e}")
        return {"message": f"failed: {e}", "page_count": 0, "file_path": ""}